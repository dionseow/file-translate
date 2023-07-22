import os
from typing import Callable
import os
from pptx import Presentation
from app.schemas.file import FileTranslator

class PowerPointTranslator(FileTranslator):
    def translate_in_place(self, file_path: str | os.PathLike, new_file_path: str | os.PathLike, translator: Callable[[str], str]) -> None:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            cell_text = cell.text
                            # Fix this
                            print(cell_text)
                            # if cell_text.strip():
                            #     translated_text = translate(text)
                            #     cell.text = cell.text.replace(cell_text, translated_text)
                                # text_frame = cell.text_frame
                                # text_frame.auto_size = True
                elif not shape.has_text_frame:
                    continue
                else:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                text = run.text
                                translated_text = translator(text)
                                run.text = run.text.replace(text, translated_text)
                    # text_frame.auto_size = True
                
        presentation.save(new_file_path)