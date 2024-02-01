import os
from typing import Callable
import os
from pptx import Presentation
from schemas.file import Processor

class PowerPointProcessor(Processor):
    def extract_text(self, file_path: str | os.PathLike) -> str:
        presentation = Presentation(file_path)
        
        extracted_text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            cell_text = cell.text
                            if cell_text.strip():
                                extracted_text = extracted_text + " " + cell_text
                
                elif not shape.has_text_frame:
                    continue

                else:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                text = run.text
                                if text.strip():
                                    extracted_text = extracted_text + " " + text
        
        return extracted_text

    def translate_in_place(self, file_path: str | os.PathLike, new_file_path: str | os.PathLike, translator: Callable[[str], str]) -> None:
        pass
    #     presentation = Presentation(file_path)
    #     for slide in presentation.slides:
    #         for shape in slide.shapes:
    #             if shape.has_table:
    #                 table = shape.table
    #                 for row in table.rows:
    #                     for cell in row.cells:
    #                         cell_text = cell.text
    #                         # Fix this
    #                         print(cell_text)
    #                         # if cell_text.strip():
    #                         #     translated_text = translate(text)
    #                         #     cell.text = cell.text.replace(cell_text, translated_text)
    #                             # text_frame = cell.text_frame
    #                             # text_frame.auto_size = True
    #             elif not shape.has_text_frame:
    #                 continue
    #             else:
    #                 text_frame = shape.text_frame
    #                 for paragraph in text_frame.paragraphs:
    #                     for run in paragraph.runs:
    #                         if run.text:
    #                             text = run.text
    #                             translated_text = translator(text)
    #                             run.text = run.text.replace(text, translated_text)
    #                 # text_frame.auto_size = True
                
    #     presentation.save(new_file_path)